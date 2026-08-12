#!/usr/bin/env python3
"""P0 KB 接口多文件类型全量端到端测试。

对 rag-engine 支持的 6 种文档格式逐一做「上传 → 解析 → 分块 → 查询」完整链路
验证，并对文件类型白名单做负向校验。

支持格式（SPEC/kb_service.proto）：
    md | txt | pdf | docx | xlsx | pptx

每个格式真实构造文件内容（二进制格式用对应库生成真实文件），全部走：
    POST /documents (GetDocumentUploadURL)
    -> PUT presigned URL 上传到 MinIO
    -> POST /notify-uploaded
    -> GET /documents 轮询 parse_status=ready
    -> POST /query 验证可检索
    -> DELETE /documents + DELETE /knowledge-bases (清理)

另含负向用例：非法 file_type（exe）应返回 400。
"""
import hashlib
import io
import json
import sys
import time
import uuid

import requests

GATEWAY = "http://localhost:8080"
TENANT_ID = "00000000-0000-0000-0000-000000000001"
USER_ID = "00000000-0000-0000-0000-000000000002"
HEADERS = {
    "X-Dev-Tenant-ID": TENANT_ID,
    "X-Dev-User-ID": USER_ID,
    "Content-Type": "application/json",
}

CHUNK_SIZE = 512


# ── 各类文件真实内容构造 ────────────────────────────────────────────────────
def _md_txt_content(kind: str) -> str:
    """md/txt：统一长文本，含 3 个章节标题，用于验证分块与检索。"""
    return (
        "# ANI 平台多文件类型解析测试（{kind}）\n\n"
        "## 第一章 作业调度\n"
        "ANI 平台的作业调度系统支持多租户、分布式与异构 GPU 的大规模调度，"
        "实现了先来先服务、优先级抢占、公平共享与基于 DRF 的占优资源公平算法。\n\n"
        "## 第二章 混合检索\n"
        "混合检索同时调用向量检索与全文检索两条通路，并通过倒数排名融合（RRF）"
        "合并排名，兼顾语义相关与字面匹配两种信号。\n\n"
        "## 第三章 对象存储\n"
        "数据管理模块基于与 S3 兼容的对象存储提供持久化能力，文件上传采用预签名"
        " URL 直传，避免大文件在网关侧产生带宽瓶颈。\n\n"
        "本文件用于验证 {kind} 格式的上传、解析、分块与检索全链路。\n"
    ).format(kind=kind)


def _build_bytes(ftype: str) -> bytes:
    """根据 file_type 构造真实文件字节内容。"""
    if ftype in ("md", "txt"):
        return _md_txt_content(ftype).encode("utf-8")
    if ftype == "pdf":
        import fitz  # PyMuPDF

        doc = fitz.open()
        text = _md_txt_content("PDF")
        for _ in range(3):
            page = doc.new_page()
            page.insert_text((72, 72), "ANI Platform PDF Parse Test")
        page = doc.new_page()
        # 大段文字以验证文本提取
        try:
            page.insert_textbox((72, 72, 540, 780), text, fontsize=10)
        except Exception:  # noqa: BLE001
            pass
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        return buf.getvalue()
    if ftype == "docx":
        import docx

        d = docx.Document()
        d.add_heading("ANI 平台多文件类型解析测试（DOCX）", level=1)
        d.add_heading("第一章 作业调度", level=2)
        d.add_paragraph(
            "ANI 平台的作业调度系统支持多租户、分布式与异构 GPU 的大规模调度，"
            "实现了先来先服务、优先级抢占、公平共享与基于 DRF 的占优资源公平算法。"
        )
        d.add_heading("第二章 混合检索", level=2)
        d.add_paragraph(
            "混合检索同时调用向量检索与全文检索两条通路，并通过倒数排名融合（RRF）"
            "合并排名，兼顾语义相关与字面匹配两种信号。"
        )
        buf = io.BytesIO()
        d.save(buf)
        return buf.getvalue()
    if ftype == "xlsx":
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "ANI"
        rows = [
            ("主题", "说明"),
            ("作业调度", "支持多租户、分布式与异构 GPU 的大规模调度"),
            ("混合检索", "向量与全文双通路，RRF 倒数排名融合"),
            ("对象存储", "S3 兼容，预签名 URL 直传"),
        ]
        for r in rows:
            ws.append(r)
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    if ftype == "pptx":
        import pptx
        from pptx.util import Inches

        prs = pptx.Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        prs.slides[0].shapes.title.text = "ANI 平台多文件类型解析测试（PPTX）"
        prs.slides[0].placeholders[1].text = (
            "作业调度支持多租户、分布式与异构 GPU 大规模调度；"
            "混合检索通过 RRF 融合向量与全文排名。"
        )
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    raise ValueError(f"unsupported ftype: {ftype}")


# ── 测试执行 ────────────────────────────────────────────────────────────────
SUMMARY = []
# 展示时保留完整值的字段（其余超长字段截断，避免刷屏）
_JSON_SAFE = {"idempotency_key", "storage_path", "upload_url", "checksum_sha256",
              "question"}


def _show(v: object, limit: int = 800) -> str:
    s = str(v)
    if len(s) <= limit:
        return s
    return s[:limit] + f"...(截断, 共{len(s)}字符)"


def _req_line(method, url, *, json_body=None, data=None, params=None):
    print(f"    REQUEST: {method} {url}")
    if params:
        print(f"      query_params = {json.dumps(params, ensure_ascii=False, default=str)}")
    if json_body is not None:
        safe = {k: (_show(v) if k in _JSON_SAFE or len(str(v)) <= 200 else _show(v, 200))
                for k, v in json_body.items()}
        print(f"      body(JSON)   = {json.dumps(safe, ensure_ascii=False, default=str, indent=2)}")
    elif data is not None:
        print(f"      body(raw)    = {_show(data)}")


def _resp_line(status, body, extra=""):
    print(f"    RESPONSE(status={status}){extra}")
    if isinstance(body, dict):
        shown = {k: (_show(v, 400) if k in _JSON_SAFE or len(str(v)) <= 400 else _show(v, 400))
                 for k, v in body.items()}
        if "sources" in body and isinstance(body["sources"], list):
            shown["sources"] = _show(body["sources"], 2000)
        if "answer" in body:
            shown["answer"] = _show(body["answer"], 1500)
        print("      body(JSON)  = " + json.dumps(shown, ensure_ascii=False, default=str, indent=2))
    else:
        print(f"      body(raw)   = {_show(body)}")


def request(method, url, *, json=None, data=None, params=None, timeout=30, put_bytes=None):
    _req_line(method, url, json_body=json, data=data, params=params)
    if put_bytes is not None:
        r = requests.put(url, data=put_bytes, timeout=timeout)
    else:
        r = requests.request(method, url, headers=HEADERS, json=json, data=data,
                             params=params, timeout=timeout)
    _resp_line(r.status_code, _safe_body(r))
    return r


def _safe_body(r):
    try:
        return r.json()
    except Exception:
        return r.text


def record(label, ok, detail=""):
    SUMMARY.append((label, ok, detail))
    print(f"  [{'PASS' if ok else 'FAIL'}] {label} {detail}")


def create_kb(ftype):
    body = {
        "idempotency_key": f"p0t_create_{ftype}_{uuid.uuid4()}",
        "name": f"p0t-{ftype}-{int(time.time())}",
        "description": f"多文件类型 P0 测试: {ftype}",
        "embedding_model": "Qwen3-Embedding-0.6B",
        "chunk_size": CHUNK_SIZE,
        "top_k": 5,
        "score_threshold": 0.3,
        "retrieval_mode": "hybrid",
    }
    r = request("POST", f"{GATEWAY}/api/v1/svc/knowledge-bases", json=body)
    kb = r.json() if r.status_code in (200, 201) else {}
    return r.status_code, kb.get("id"), r


def upload_doc(kb_id, ftype, content_bytes):
    req = {
        "idempotency_key": f"p0t_upload_{ftype}_{uuid.uuid4()}",
        "file_name": f"sample.{ftype}",
        "file_type": ftype,
        "file_size_bytes": len(content_bytes),
        "checksum_sha256": hashlib.sha256(content_bytes).hexdigest(),
    }
    r = request("POST", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents", json=req)
    body = r.json() if r.status_code == 200 else {}
    # 真实上传到 MinIO（预签名 URL 直传）
    up = request("PUT", body.get("upload_url") or "", put_bytes=content_bytes)
    ntf = {}
    if up.status_code in (200, 204):
        nn = request("POST",
                     f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents/notify-uploaded",
                     json={"doc_id": body.get("doc_id"), "storage_path": body.get("storage_path")})
        ntf = {"notify_status": nn.status_code}
    return r, body, up, ntf


def wait_ready(kb_id, doc_id, timeout=180):
    t0 = time.time()
    last = None
    while time.time() - t0 < timeout:
        r = request("GET", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents",
                    params={"limit": "20"})
        if r.status_code == 200:
            for it in r.json().get("items", []):
                if it.get("doc_id") == doc_id or it.get("id") == doc_id:
                    last = it
                    st = it.get("parse_status")
                    print(f"      parse_status={st} chunk_count={it.get('chunk_count')}")
                    if st in ("ready", "failed"):
                        return st, it.get("chunk_count") or 0
        time.sleep(5)
    return (last or {}).get("parse_status"), (last or {}).get("chunk_count") or 0


def run_query(kb_id):
    req = {
        "idempotency_key": f"p0t_query_{uuid.uuid4()}",
        "question": "ANI 平台的作业调度能力与混合检索原理是什么？",
        "top_k": 5,
        "score_threshold": 0.3,
        "retrieval_mode": "hybrid",
    }
    r = request("POST", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/query",
                json=req, timeout=120)
    raw = r.json() if r.status_code == 200 else {}
    srcs = raw.get("sources") if isinstance(raw, dict) and raw.get("sources") else []
    ok = r.status_code == 200 and len(srcs) > 0 and all((s.get("content") or "").strip() for s in srcs)
    return ok, r.status_code, len(srcs), raw


def main():
    print("=" * 78)
    print("  P0 多文件类型全链路测试 (chunk_size=%d)" % CHUNK_SIZE)
    print("=" * 78)

    ftypes = ["md", "txt", "pdf", "docx", "xlsx", "pptx"]

    for ftype in ftypes:
        sep = "=" * 78
        print(f"\n{sep}\n  ▶ 文件类型: {ftype}\n{sep}")
        content_bytes = _build_bytes(ftype)
        print(f"      构造文件字节 = {len(content_bytes)} ({ftype})")

        # 1) 建 KB
        st, kb_id, _ = create_kb(ftype)
        print(f"      CreateKB status={st} kb_id={kb_id}")
        record(f"[{ftype}] CreateKB", st == 201, f"kb_id={kb_id}")
        if st != 201:
            continue

        # 2) 上传 + 真实 PUT + notify
        r, body, up, ntf = upload_doc(kb_id, ftype, content_bytes)
        doc_id = body.get("doc_id")
        print(f"      GetDocumentUploadURL status={r.status_code} doc_id={doc_id}")
        print(f"      MinIO PUT status={up.status_code}  notify status={ntf.get('notify_status')}")
        record(f"[{ftype}] GetDocumentUploadURL", r.status_code == 200 and doc_id)
        record(f"[{ftype}] MinIO PUT", up.status_code in (200, 204))
        record(f"[{ftype}] NotifyUploaded", ntf.get("notify_status") == 202)

        # 3) 轮询解析
        pstatus, chunk_count = wait_ready(kb_id, doc_id)
        print(f"      parse_status={pstatus} chunk_count={chunk_count}")
        record(f"[{ftype}] 解析==ready", pstatus == "ready", f"status={pstatus}")
        record(f"[{ftype}] chunk_count>0", chunk_count > 0, f"chunk_count={chunk_count}")
        if pstatus != "ready":
            # 解析失败也继续清理，不再查询
            request("DELETE", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents/{doc_id}")
            request("DELETE", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}")
            continue

        # 4) 查询验证可检索
        q_ok, q_st, ns, raw = run_query(kb_id)
        print(f"      Query status={q_st} sources={ns}")
        record(f"[{ftype}] Query 可检索", q_ok, f"status={q_st} sources={ns}")
        if q_ok:
            first = raw["sources"][0]
            print(f"        top source file_name={first.get('file_name')} score={first.get('score')}")

        # 5) 清理
        d1 = request("DELETE", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents/{doc_id}")
        d2 = request("DELETE", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}")
        record(f"[{ftype}] DeleteDocument", d1.status_code == 204)
        record(f"[{ftype}] DeleteKB", d2.status_code == 204)

    # ── 负向：非法文件类型 ─────────────────────────────────────────────────
    print(f"\n{'=' * 78}\n  ▶ 负向用例: 非法 file_type 白名单校验\n{'=' * 78}")
    st_bad, kb_id2, _ = create_kb("md")
    if st_bad == 201:
        req = {
            "idempotency_key": f"p0t_bad_{uuid.uuid4()}",
            "file_name": "evil.exe",
            "file_type": "exe",
            "file_size_bytes": 10,
            "checksum_sha256": "0" * 64,
        }
        r = request("POST", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id2}/documents", json=req)
        msg = r.json().get("message", "") if r.headers.get("content-type", "").startswith("application/json") else r.text
        print(f"      upload file_type=exe status={r.status_code}")
        record("非法 file_type(exe) 被拒", r.status_code == 400, f"status={r.status_code} msg={msg}")
        request("DELETE", f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id2}")
    else:
        record("非法 file_type(exe) 被拒", False, f"前置 CreateKB 失败 status={st_bad}")

    # ── 汇总 ───────────────────────────────────────────────────────────────
    print(f"\n{'=' * 78}\n  结果汇总\n{'=' * 78}")
    passed = sum(1 for _, ok, _ in SUMMARY if ok)
    failed = len(SUMMARY) - passed
    print(f"  通过: {passed}")
    print(f"  失败: {failed}")
    for label, ok, detail in SUMMARY:
        print(f"    [{'PASS' if ok else 'FAIL'}] {label} {detail}")
    return 0 if failed == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
