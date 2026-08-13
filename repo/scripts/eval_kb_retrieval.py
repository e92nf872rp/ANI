#!/usr/bin/env python3
"""多文件类型 × 三种检索模式 的检索质量（Precision/Recall/MRR）评估。

重写版说明（相对旧 eval_kb_retrieval.py 的修正）：
  1. 指标口径修正：
     - Precision@K = 每个查询在“返回的前 K 个 source（不足则按实际返回数）”中
       相关命中数 的占比，再对同模式下的所有查询做宏平均。
     - Recall@K   = 每个查询“在返回的前 K 个中是否召回到该主题”（0/1），再宏平均。
     - MRR        = 首个相关 source 的倒数排名，无命中记 0。
     修正了旧脚本把全部返回累计计数、导致 Precision@K 与 Recall@K 错乱
     （甚至 Recall>1）的问题。
  2. 相关判定：用每个主题的“强标识词”（专有机制名，如 DRF / RRF / 预签名URL /
     HNSW / 行级安全 / 阈值兜底），判定返回 source 是否确属该主题，避免把
     doc_summary 等泛述误判为相关。
  3. 覆盖 6 种文件类型（md/txt/pdf/docx/xlsx/pptx）× 3 种检索模式
     (hybrid/vector/keyword)。

指标定义（对每个查询 q 返回的 source 序列 S）：
    rel(i) = 1 若第 i 个 source 命中该主题的强标识词，否则 0
    Precision@K(q) = sum_{i<K} rel(i) / min(K, len(S))
    Recall@K(q)    = 1 if sum_{i<K} rel(i) >= 1 else 0
    MRR(q)         = 1 / (第一个 rel 所在位次+1)，无 rel 记 0
宏平均在 mode / ftype 维度给出。
"""
import argparse
import hashlib
import io
import json
import sys
import time
import uuid

import requests

GATEWAY = "http://localhost:8080"
TENANT_ID = "00000000-0000-0000-0000-000000000001"
USER_ID = "00000000-0000-0000-0000-000000000002"
HEADERS = {
    "X-Dev-Tenant-ID": TENANT_ID,
    "X-Dev-User-ID": USER_ID,
    "Content-Type": "application/json",
}

CHUNK_SIZE = 512
MODES = ["hybrid", "vector", "keyword"]
K_VALUES = [1, 3, 5]
FTYPES = ["md", "txt", "pdf", "docx", "xlsx", "pptx"]

# 每个主题：主题名 → (唯一强标识词集, 查询问题)。标识词用于判定 source 是否相关。
TOPICS = {
    "作业调度": (
        ("DRF", "占优资源公平", "先来先服务", "优先级抢占"),
        "ANI 平台在大规模训练任务之间如何公平地分配 GPU 资源？",
    ),
    "混合检索": (
        ("RRF", "倒数排名融合", "混合检索", "向量检索与全文检索"),
        "平台把向量检索与全文检索两条通路的结果合并排序时使用哪种算法？",
    ),
    "对象存储": (
        ("预签名URL", "预签名", "S3", "SHA-256"),
        "上传大文件时平台通过什么方式直传对象存储从而避免网关产生带宽瓶颈？",
    ),
}


def build_docs_text(kind: str) -> str:
    """构造含 3 个主题的 markdown 文本（每个主题一个独立章节）。"""
    header = f"# ANI 平台检索质量评估手册（{kind}）\n\n本手册描述 ANI 平台的多项核心能力，用于评估不同检索模式的命中效果。\n"
    sections = [
        (
            "## 作业调度\n"
            "ANI 平台的作业调度系统面向 AI 训练与推理提供大规模多租户调度能力。"
            "在多租户同时提交海量任务的场景下，系统采用 DRF 占优资源公平算法，"
            "以最小共享量的占优份额作为衡量标准，在各租户与各作业之间按占优资源公平分配，"
            "既能保证资源公平又能有效利用异构 GPU 集群的空闲算力。调度器支持先来先服务、"
            "优先级抢占与公平共享等多种策略，可与 DRF 占优资源公平算法配合使用，"
            "并支持多机多卡的一体化分配。\n"
        ),
        (
            "## 混合检索\n"
            "平台提供一种融合语义信号与字面信号的混合检索方案。它同时调用向量检索与全文检索"
            "两条通路，并把两条通路的结果通过倒数排名融合（RRF）算法合并排序："
            "根据每个分块在两条通路中的排名计算融合得分，排名越靠前累计加分越高，"
            "从而让向量分与文本分这两个量纲不同的分数能够公平参与最终排序。"
            "该机制兼顾语义相关与关键词匹配两种信号，显著提升复杂查询的召回效果。\n"
        ),
        (
            "## 对象存储\n"
            "平台数据层基于与 S3 兼容的对象存储提供持久化能力。文件上传采用预签名URL机制："
            "网关先行签发带时效与权限范围的凭证与上传地址，客户端随后直接向对象存储发起上传请求，"
            "大文件数据无需在网关侧中转，从机制上避免网关成为带宽瓶颈，降低单点开销并提升吞吐。"
            "同时为每个上传文档计算 SHA-256 校验和以保证内容完整性。\n"
        ),
    ]
    return header + "\n".join(sections)


def build_bytes(ftype: str, kind: str) -> bytes:
    """按 file_type 构造真实文件字节。"""
    text = build_docs_text(kind)
    if ftype in ("md", "txt"):
        return text.encode("utf-8")
    if ftype == "pdf":
        import fitz

        doc = fitz.open()
        page = doc.new_page()
        try:
            page.insert_textbox((72, 72, 540, 780), text, fontsize=9)
        except Exception:  # noqa: BLE001
            # fallback 逐行写入
            y = 72
            for line in text.splitlines():
                try:
                    page.insert_text((72, y), line[:80], fontsize=9)
                except Exception:  # noqa: BLE001
                    pass
                y += 14
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        return buf.getvalue()
    if ftype == "docx":
        import docx

        d = docx.Document()
        d.add_heading(f"ANI 平台检索质量评估手册（DOCX）", level=1)
        for sec in ("## 作业调度\n", "## 混合检索\n", "## 对象存储\n"):
            d.add_paragraph(sec.replace("## ", ""))
            if sec.startswith("## 作业"):
                d.add_paragraph(
                    "调度系统采用 DRF 占优资源公平算法，以最小共享占优份额公平分配 GPU 资源，"
                    "支持先来先服务、优先级抢占与公平共享策略，兼容多机多卡一体化调度。"
                )
            elif sec.startswith("## 混合"):
                d.add_paragraph(
                    "混合检索同时调用向量检索与全文检索，并以倒数排名融合 RRF 合并排序，"
                    "兼顾语义相关与关键词匹配信号。"
                )
            else:
                d.add_paragraph(
                    "对象存储采用 S3 兼容协议，文件上传使用预签名URL 直传，"
                    "为避免网关带宽瓶颈计算 SHA-256 校验和保证完整性。"
                )
        buf = io.BytesIO()
        d.save(buf)
        return buf.getvalue()
    if ftype == "xlsx":
        import openpyxl

        wb = openpyxl.Workbook()
        for title, tagline in (
            ("作业调度", "DRF 占优资源公平算法；先来先服务；优先级抢占；多机多卡调度"),
            ("混合检索", "RRF 倒数排名融合；向量检索与全文检索双通路"),
            ("对象存储", "S3 兼容；预签名URL 直传；SHA-256 校验和"),
        ):
            ws = wb.create_sheet(title)
            ws.append(("主题", "说明"))
            ws.append((title, tagline))
        if "Sheet" in wb.sheetnames:
            del wb["Sheet"]
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    if ftype == "pptx":
        import pptx
        from pptx.util import Inches

        prs = pptx.Presentation()
        for title, body in (
            ("作业调度", "DRF 占优资源公平算法，支持先来先服务、优先级抢占、公平共享与多机多卡一体化调度。"),
            ("混合检索", "RRF 倒数排名融合算法，合并向量检索与全文检索两条通路的排名。"),
            ("对象存储", "S3 兼容对象存储，预签名URL 直传，SHA-256 校验和保证完整性。"),
        ):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    raise ValueError(f"unsupported ftype: {ftype}")


def is_hit(content: str, keywords: tuple) -> bool:
    """返回 source 是否属于该主题：命中任一强标识词。"""
    if not content:
        return False
    return any(k in content for k in keywords)


def compute_metrics(sources, keywords):
    """按返回 source 序列计算各 K 的 Precision/Recall 与 MRR。"""
    n = len(sources)
    rel = [1 if is_hit(s.get("content") or "", keywords) else 0 for s in sources]
    first = next((i for i, v in enumerate(rel) if v), None)
    mrr = 1.0 / (first + 1) if first is not None else 0.0
    prec, rec = {}, {}
    for K in K_VALUES:
        window = rel[:K]
        denom = min(K, n)
        prec[K] = sum(window) / denom if denom else 0.0
        rec[K] = 1.0 if sum(window) >= 1 else 0.0
    return prec, rec, mrr


# ── HTTP 帮助 ─────────────────────────────────────────────────────────────
def _safe_body(r):
    try:
        return r.json()
    except Exception:  # noqa: BLE001
        return r.text


def create_kb(ftype):
    body = {
        "idempotency_key": f"eval_create_{ftype}_{uuid.uuid4()}",
        "name": f"eval-{ftype}-{int(time.time())}",
        "description": f"检索质量评估: {ftype}",
        "embedding_model": "Qwen3-Embedding-0.6B",
        "chunk_size": CHUNK_SIZE,
        "top_k": 5,
        "score_threshold": 0.3,
        "retrieval_mode": "hybrid",
    }
    r = requests.post(f"{GATEWAY}/api/v1/svc/knowledge-bases", headers=HEADERS, json=body, timeout=30)
    try:
        kb = r.json()
    except Exception:  # noqa: BLE001
        kb = {}
    return r.status_code, kb.get("id")


def upload_doc(kb_id, ftype, content_bytes):
    req = {
        "idempotency_key": f"eval_upload_{ftype}_{uuid.uuid4()}",
        "file_name": f"sample.{ftype}",
        "file_type": ftype,
        "file_size_bytes": len(content_bytes),
        "checksum_sha256": hashlib.sha256(content_bytes).hexdigest(),
    }
    r = requests.post(f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents",
                      headers=HEADERS, json=req, timeout=30)
    body = r.json() if r.status_code == 200 else {}
    up = requests.put(body.get("upload_url") or "", data=content_bytes, timeout=60)
    ntf = 0
    if up.status_code in (200, 204) and body.get("doc_id"):
        nn = requests.post(
            f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents/{body['doc_id']}/notify-uploaded",
            headers=HEADERS, json={"storage_path": body.get("storage_path")}, timeout=30)
        ntf = nn.status_code
    return r.status_code, body.get("doc_id"), up.status_code, ntf


def wait_ready(kb_id, doc_id, timeout=240):
    t0 = time.time()
    while time.time() - t0 < timeout:
        r = requests.get(f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents",
                         headers=HEADERS, params={"limit": "20"}, timeout=30)
        if r.status_code == 200:
            for it in r.json().get("items", []):
                if it.get("doc_id") == doc_id or it.get("id") == doc_id:
                    st = it.get("parse_status")
                    if st in ("ready", "failed"):
                        return st, it.get("chunk_count") or 0
        time.sleep(5)
    return "timeout", 0


def run_query(kb_id, question, mode):
    req = {
        "idempotency_key": f"eval_query_{mode}_{uuid.uuid4()}",
        "question": question,
        "top_k": 5,
        "score_threshold": 0.3,
        "retrieval_mode": mode,
    }
    r = requests.post(f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/query",
                      headers=HEADERS, json=req, timeout=180)
    raw = r.json() if r.status_code == 200 else {}
    sources = raw.get("sources") if isinstance(raw, dict) else []
    return r.status_code, sources


def cleanup(kb_id, doc_id):
    if doc_id:
        requests.delete(f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}/documents/{doc_id}",
                        headers=HEADERS, timeout=30)
    requests.delete(f"{GATEWAY}/api/v1/svc/knowledge-bases/{kb_id}", headers=HEADERS, timeout=30)


def main():
    ap = argparse.ArgumentParser(description="KB 检索质量评估（多文件类型 × 三种模式）")
    ap.add_argument("--ftypes", nargs="*", default=FTYPES, help="文件类型子集")
    ap.add_argument("--modes", nargs="*", default=MODES, help="检索模式子集")
    ap.add_argument("--topics", nargs="*", default=list(TOPICS.keys()), help="主题子集")
    ap.add_argument("--pdf", action="store_true", help="(已默认) 忽略")
    args = ap.parse_args()

    ftypes = args.ftypes
    modes = args.modes
    topic_keys = [t for t in args.topics if t in TOPICS]

    print("=" * 80)
    print(f"  KB 检索质量评估  chunk_size={CHUNK_SIZE}  模式={modes}")
    print(f"  文件类型: {ftypes}  主题: {topic_keys}")
    print("=" * 80)

    # results[ftype][mode] = list[dict] （每个 query 一条）
    results = {f: {m: [] for m in modes} for f in ftypes}

    for ftype in ftypes:
        print(f"\n{'─' * 80}\n  ▶ 文件类型: {ftype}\n{'─' * 80}")
        content_bytes = build_bytes(ftype, ftype.upper())
        print(f"  构造文件字节 = {len(content_bytes)} ({ftype})")

        st, kb_id = create_kb(ftype)
        print(f"  CreateKB status={st} kb_id={kb_id}")
        if st != 201 or not kb_id:
            for m in modes:
                results[ftype][m].append({"error": f"CreateKB {st}"})
            continue

        st2, doc_id, up_st, ntf_st = upload_doc(kb_id, ftype, content_bytes)
        print(f"  Upload status={st2} doc_id={doc_id} PUT={up_st} notify={ntf_st}")
        if st2 != 200 or not doc_id:
            cleanup(kb_id, None)
            continue

        pstatus, cc = wait_ready(kb_id, doc_id)
        print(f"  parse_status={pstatus} chunk_count={cc}")
        if pstatus != "ready":
            print(f"  [SKIP] {ftype} 解析未 ready，跳过检索评估")
            cleanup(kb_id, doc_id)
            continue

        for topic in topic_keys:
            keywords, question = TOPICS[topic]
            for mode in modes:
                qst, sources = run_query(kb_id, question, mode)
                row = {"topic": topic, "question": question, "status": qst,
                       "returned": len(sources)}
                if qst == 200:
                    prec, rec, mrr = compute_metrics(sources, keywords)
                    row.update({f"p@{K}": prec[K] for K in K_VALUES})
                    row.update({f"r@{K}": rec[K] for K in K_VALUES})
                    row["mrr"] = mrr
                    row["top_src"] = (sources[0].get("content") or "")[:40].replace("\n", " ") if sources else ""
                else:
                    for K in K_VALUES:
                        row[f"p@{K}"] = 0.0
                        row[f"r@{K}"] = 0.0
                    row["mrr"] = 0.0
                results[ftype][mode].append(row)
                print(f"    [{ftype}][{mode}][{topic}] status={qst} returned={len(sources)} "
                      f"P@3={row.get('p@3'):.3f} R@3={row.get('r@3'):.3f} MRR={row.get('mrr'):.3f}")

        cleanup(kb_id, doc_id)

    # ── 汇总 ─────────────────────────────────────────────────────────────
    print("\n" + "=" * 80)
    print("  汇总：各文件类型 × 各检索模式 的宏平均指标")
    print("=" * 80)

    avg_keys = [f"p@{K}" for K in K_VALUES] + [f"r@{K}" for K in K_VALUES] + ["mrr"]
    print("\n[按文件类型 × 模式]")
    header = f"{'ftype':<7}{'mode':<9}" + "".join(f"{k:>9}" for k in avg_keys)
    print(header)
    mode_agg = {m: {k: [] for k in avg_keys} for m in modes}
    for ft in ftypes:
        for m in modes:
            rows = [r for r in results[ft][m] if "error" not in r]
            if not rows:
                print(f"{ft:<7}{m:<9}" + "".join(f"{'-':>9}" for _ in avg_keys))
                continue
            vals = {}
            for k in avg_keys:
                vals[k] = sum(r.get(k, 0.0) for r in rows) / len(rows)
                mode_agg[m][k].append(vals[k])
            line = f"{ft:<7}{m:<9}" + "".join(f"{vals[k]:>9.3f}" for k in avg_keys)
            print(line)

    print("\n[按检索模式 全局宏平均]")
    print(f"{'mode':<9}" + "".join(f"{k:>9}" for k in avg_keys))
    for m in modes:
        if not mode_agg[m][avg_keys[0]]:
            print(f"{m:<9}" + "".join(f"{'-':>9}" for _ in avg_keys))
            continue
        line = f"{m:<9}" + "".join(f"{sum(mode_agg[m][k])/len(mode_agg[m][k]):>9.3f}" for k in avg_keys)
        print(line)

    # 摘要
    print("\n[摘要] 三种模式全局平均")
    for m in modes:
        if not mode_agg[m][avg_keys[0]]:
            continue
        p5 = sum(mode_agg[m]["p@5"]) / len(mode_agg[m]["p@5"])
        r5 = sum(mode_agg[m]["r@5"]) / len(mode_agg[m]["r@5"])
        mrr = sum(mode_agg[m]["mrr"]) / len(mode_agg[m]["mrr"])
        f1 = 2 * p5 * r5 / (p5 + r5) if (p5 + r5) else 0.0
        print(f"  {m:<9} P@5={p5:.3f}  R@5={r5:.3f}  MRR={mrr:.3f}  F1@5={f1:.3f}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
